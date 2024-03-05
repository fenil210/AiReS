import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates2 import new, bot_template, user_template,page_bg_img
from langchain_community.llms import CTransformers
from langchain.llms import HuggingFaceHub
import time
from pptx import Presentation
from docx import Document
import io 
import os

def pdfread(pdf):
        reader=PdfReader(pdf)
        text=[]
        for i in range(len(reader.pages)):
            text.append(reader.pages[i].extract_text())

        return text

def readlat(lat):
  te=LatexNodes2Text().latex_to_text(lat)
  te.strip()
  return te

def readocx(poc):
  poi=io.BytesIO(poc)
  document = Document(poi)
  text=[]
  for x in (document.paragraphs):
    text.append(x.text)
  return text

def readppt(pp):
  pio=io.BytesIO(pp)
  prs=Presentation(pio)
  text=[]
  for slide in prs.slides:
    temp=[]
    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      for p in shape.text_frame.paragraphs:
        for r in p.runs:
          temp.append(r.text)
      for x in temp:
        text.append(x)
  return text

def typewriter(text, template, speed):
    tokens = text.split()
    container = st.empty()
    for index in range(len(tokens) + 1):
        curr_full_text = " ".join(tokens[:index])
        container.markdown(curr_full_text)
        time.sleep(1 / speed)

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    # llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})
   
                       
    
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    retrieved_documents = conversation_chain.run(query="your_query")

    # Loop through each retrieved document and access its metadata
    for doc in retrieved_documents:
        print(f"Document Text: {doc.text}")
        print(f"Document Metadata: {doc.metadata}")
    
    return conversation_chain


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']
    
    mes = st.session_state.chat_history[-1]
    if len(st.session_state.chat_history)%2==1:
        typewriter(mes.content, user_template, 10)
    else:
        typewriter(mes.content, bot_template, 10)

def submit():
    st.session_state.my_text=st.session_state.widget
    st.session_state.widget=""


def main():
    load_dotenv()
    st.set_page_config(page_title="Ai4Researchers",page_icon=":cactus:")
    st.markdown(page_bg_img,unsafe_allow_html=True)
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None
    st.title("Ai4Researchers")
    st.header("Introducing")
    st.write(new,unsafe_allow_html=True)
    if "my_text" not in st.session_state:
        st.session_state.my_text=""
    st.text_input("Hello! This is your research Assistant",key="widget",on_change=submit)
    user_question=st.session_state.my_text
    if user_question:
        st.subheader(user_question)
        handle_userinput(user_question)

    st.subheader("Document Library")
    docs = st.file_uploader("Upload your PDFs here and click on 'Process'", accept_multiple_files=True)
    if st.button("Process"):
        with st.spinner("Processing"):
            # get pdf text
            raw_text = ""
            print(docs)
            for x in docs:
                ty=(x.name).split('.')
                if('pdf'==ty[1]):
                    te=pdfread(x)
                    raw_text+=' '.join(te)
                elif('docx'==ty[1]):
                    te=readocx(x.read())
                    raw_text+=' '.join(te)
                elif('pptx'==ty[1]):
                    te=readppt(x.read())
                    raw_text+=' '.join(te)
                elif('tex'==ty[1]):
                    bd=x.getvalue()
                    so=StringIO(bd.decode("utf-8"))
                    so=so.read()
                    raw_text+=readlat(so)
            text_chunks = get_text_chunks(raw_text)
            vectorstore = get_vectorstore(text_chunks)
            st.session_state.conversation = get_conversation_chain(
                    vectorstore)


if __name__ == '__main__':
    main()
